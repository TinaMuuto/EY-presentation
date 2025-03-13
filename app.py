import streamlit as st
import pandas as pd
import pptx
import io
import copy
import requests
import re
from PIL import Image
from pptx.dml.color import RGBColor
import time

# Undgå advarsler ved store billeder
Image.MAX_IMAGE_PIXELS = None

# Progress bar
steps = ["Upload fil", "Behandler data", "Genererer slides", "Færdig - Klar til download"]
progress = st.progress(0)
status_text = st.empty()

def update_progress(step):
    """Opdaterer progress bar og status-tekst"""
    progress.progress((step + 1) / len(steps))
    status_text.text(steps[step])
    time.sleep(1)

#############################################
# Hjælpefunktioner – Dataopslag
#############################################

def normalize_variantkey(s):
    return str(s).lower().replace(" - config", "").strip()

def match_variant_rows(item_no, mapping_df):
    item_no_norm = normalize_variantkey(item_no)
    mapping_df.columns = mapping_df.columns.str.strip().str.replace(r"[{}]", "", regex=True)
    
    # Log kolonnenavne for fejlfinding
    st.write("Kolonner i mapping_df:", mapping_df.columns.tolist())
    
    if "Product code" in mapping_df.columns:
        mapping_df['VariantKey_norm'] = mapping_df['Product code'].apply(normalize_variantkey)
    elif "{{Product code}}" in mapping_df.columns:
        mapping_df['VariantKey_norm'] = mapping_df['{{Product code}}'].apply(normalize_variantkey)
    else:
        st.error("Fejl: 'Product code' ikke fundet i mapping-filen. Kontroller filens format.")
        st.stop()
    
    matches = mapping_df[mapping_df['VariantKey_norm'] == item_no_norm]
    if not matches.empty:
        return matches
    if '-' in item_no_norm:
        part = item_no_norm.split('-')[0]
        matches = mapping_df[mapping_df['VariantKey_norm'].str.startswith(part)]
        if not matches.empty:
            return matches
    return pd.DataFrame()

def lookup_single_value(item_no, mapping_df, column_name):
    rows = match_variant_rows(item_no, mapping_df)
    if rows.empty:
        return "N/A"
    val = rows.iloc[0].get(column_name, "N/A")
    return "N/A" if pd.isna(val) else str(val).strip()

#############################################
# Slide-håndtering – Bevar layout og design
#############################################

def copy_slide_from_template(template_slide, target_pres):
    try:
        blank_layout = target_pres.slide_layouts[6]
    except IndexError:
        blank_layout = target_pres.slide_layouts[0]
    new_slide = target_pres.slides.add_slide(blank_layout)
    for shape in template_slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide

def replace_text_placeholders(slide, replacements):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                new_text = full_text
                for key, val in replacements.items():
                    placeholder = f"{{{{{key}}}}}"
                    new_text = new_text.replace(placeholder, val)
                if new_text != full_text:
                    paragraph.clear()
                    paragraph.add_run().text = new_text

#############################################
# Billedindsættelse – Nedskaler store billeder
#############################################

def insert_image(slide, placeholder, image_url):
    if not image_url:
        return
    try:
        resample_filter = Image.LANCZOS
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip() == f"{{{{{placeholder}}}}}":
                left, top = shape.left, shape.top
                max_w, max_h = shape.width, shape.height
                resp = requests.get(image_url, timeout=10)
                if resp.status_code == 200:
                    img_data = io.BytesIO(resp.content)
                    with Image.open(img_data) as im:
                        MAX_SIZE = (1920, 1080)
                        im.thumbnail(MAX_SIZE, resample=resample_filter)
                        orig_w, orig_h = im.size
                        scale = min(max_w / orig_w, max_h / orig_h)
                        new_w = int(orig_w * scale)
                        new_h = int(orig_h * scale)
                        resized_im = im.resize((new_w, new_h), resample=resample_filter)
                        output_io = io.BytesIO()
                        resized_im.save(output_io, format="PNG")
                        output_io.seek(0)
                    slide.shapes.add_picture(output_io, left, top, width=new_w, height=new_h)
                    shape.text = ""
    except Exception as e:
        st.warning(f"Kunne ikke hente billede fra {image_url}: {e}")

#############################################
# Udfyld slides med tekst og billeder
#############################################

def fill_slide(slide, product_row, mapping_df):
    item_no = str(product_row.get("Item no", "")).strip()
    replacements = {
        "Product name": f"Product Name: {lookup_single_value(item_no, mapping_df, 'Product name')}",
        "Product code": f"Product Code: {item_no}",
        "Product country of origin": f"Country of Origin: {lookup_single_value(item_no, mapping_df, 'Country of origin')}",
    }
    replace_text_placeholders(slide, replacements)
    packshot_url = lookup_single_value(item_no, mapping_df, "Packshot image")
    insert_image(slide, "Product Packshot1", packshot_url)

#############################################
# Streamlit UI
#############################################

st.title("PowerPoint Generator")
user_file = st.file_uploader("Upload din Excel-fil", type=["xlsx"])
if user_file:
    update_progress(1)
    user_df = pd.read_excel(user_file)
    mapping_df = pd.read_excel("mapping-file.xlsx")
    prs = pptx.Presentation("template-generator.pptx")
    update_progress(2)
    
    if not prs.slides:
        st.error("Template-præsentationen har ingen slides.")
        st.stop()
    
    template_slide = prs.slides[0]
    for idx, row in user_df.iterrows():
        if idx == 0:
            slide = prs.slides[0]
        else:
            slide = copy_slide_from_template(template_slide, prs)
        fill_slide(slide, row, mapping_df)
    
    update_progress(3)
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    st.success("Præsentationen er færdig!")
    st.download_button("Download præsentation", data=ppt_io, file_name="presentation.pptx")
